"""
create_pdf_video.py
===================
直接用 Startup_Profit_Blueprint.pdf 的精美投影片製作影音：
  1. 將 PDF 每頁渲染為 1920×1080 PNG
  2. 用 gTTS 生成創投家風格中文語音旁白
  3. 用 ffmpeg 組合成 MP4
"""
import os, subprocess, glob, json
import fitz          # pymupdf
import imageio_ffmpeg

PDF_SRC  = r"C:\Users\User\Downloads\Startup_Profit_Blueprint.pdf"
BASE     = r"D:\wi"
ASSETS   = os.path.join(BASE, "pres_pdf_assets")
OUT_PPTX = os.path.join(BASE, "hw6_presentation_pdf.pptx")
OUT_MP4  = os.path.join(BASE, "hw6_presentation_pdf.mp4")
FFMPEG   = imageio_ffmpeg.get_ffmpeg_exe()

os.makedirs(ASSETS, exist_ok=True)

# ── 1. PDF 頁面 → PNG（1920×1080）─────────────────────────────────────────────
TARGET_W, TARGET_H = 1920, 1080

print("=" * 60)
print("步驟 1：將 PDF 轉換為 PNG 投影片")
print("=" * 60)

doc = fitz.open(PDF_SRC)
n_pages = len(doc)
print(f"PDF 共 {n_pages} 頁")

slide_pngs = []
for i in range(n_pages):
    page = doc[i]
    # 計算縮放比例以達到目標解析度
    rect = page.rect
    scale_x = TARGET_W / rect.width
    scale_y = TARGET_H / rect.height
    scale   = min(scale_x, scale_y)   # 等比縮放

    mat = fitz.Matrix(scale, scale)
    pix = page.get_pixmap(matrix=mat, alpha=False)

    # 如果不是剛好 1920×1080，貼到白底畫布
    import PIL.Image, io
    img_bytes = pix.tobytes("png")
    img = PIL.Image.open(io.BytesIO(img_bytes)).convert("RGB")

    canvas = PIL.Image.new("RGB", (TARGET_W, TARGET_H), (255, 255, 255))
    offset_x = (TARGET_W - img.width)  // 2
    offset_y = (TARGET_H - img.height) // 2
    canvas.paste(img, (offset_x, offset_y))

    png_path = os.path.join(ASSETS, f"slide_{i:02d}.png")
    canvas.save(png_path, "PNG")
    slide_pngs.append(png_path)
    print(f"  頁面 {i+1:02d}/{n_pages} → {os.path.basename(png_path)}  ({img.width}×{img.height})")

doc.close()

# ── 2. 創投家風格旁白腳本（12 頁）─────────────────────────────────────────────
NARRATIONS = [
    # 0 封面
    (
        "各位投資人，歡迎。"
        "今天我們要揭開一個重要真相：新創公司的獲利能力，是可以被精準預測的。"
        "這份報告基於 CRISP-DM 六階段機器學習管線，"
        "對五十家新創企業進行多元線性迴歸分析，"
        "提供科學化的資本配置藍圖。讓我們開始。"
    ),
    # 1 投資者的核心叩問
    (
        "每位投資人都在問同一個問題："
        "給定一家新創的資金分配，我應該期待多少回報？"
        "我們把這個問題轉化為機器學習任務。"
        "輸入：研發支出、行政管理、行銷支出、公司所在州。"
        "輸出：年度利潤。"
        "目標：測試集 R 平方達到零點九零——這是我們對自己設下的高標準。"
    ),
    # 2 系統健康度儀表板
    (
        "在建模之前，我們先做了完整的資料健康檢查。"
        "零值分析——保留，因為零研發支出反映真實的商業狀態。"
        "偏態分析——無需對數轉換。"
        "異常值——僅一筆邊緣異常，無需移除。"
        "共線性——VIF 全部低於五，安全。"
        "但有兩個警示燈："
        "研發與利潤相關係數高達零點九七，信號強到需要特別處理；"
        "特徵尺度差異巨大，必須標準化。"
        "資料健康，才能建出值得信任的模型。"
    ),
    # 3 尺度失衡
    (
        "看這個對比。"
        "原始資料中，行銷支出最高達四十七萬美元，研發支出最高只有十六萬。"
        "如果直接建模，演算法會誤以為數值大的行銷支出最重要——但事實正好相反。"
        "StandardScaler 標準化之後，三個特徵的均值歸零，標準差統一為一。"
        "這才是公平競爭的起跑線。"
        "現在，讓特徵的真實預測力說話。"
    ),
    # 4 篩選熔爐
    (
        "五種特徵選擇演算法同時上場投票。"
        "研發支出——五票全中，保留無懸念。"
        "行銷支出——四票，保留。"
        "行政管理——三票，存在爭議但保留。"
        "州別變數——零票，一票都沒有，一致出局。"
        "地理位置被所有演算法拋棄。"
        "這個多數決機制，成功過濾無效雜訊，讓真正的財務信號脫穎而出。"
    ),
    # 5 多次模型效能對決
    (
        "最讓人驚訝的發現在這裡。"
        "三種特徵組合對決：Correlation 方法保留研發加行銷，R 平方零點九一六八。"
        "共識方法保留研發加行政加行銷，R 平方零點九零零一。"
        "而 Tree 模型只保留一個變數——研發支出——"
        "R 平方達到零點九二六五，誤差只有七千七百一十四美元。"
        "單一特徵，最佳表現。"
        "這是研發支出作為獲利核心引擎的最直接證明。"
    ),
    # 6 管線架構
    (
        "我們的模型架構使用 Pipeline 物件，把標準化與建模封裝在一起。"
        "關鍵原因：防止資料洩漏。"
        "StandardScaler 只在訓練集配適，"
        "交叉驗證和測試集絕對不會偷看到未來的數據。"
        "這是機器學習工程的基本紀律，"
        "也是讓我們的結果值得信任的根本保障。"
    ),
    # 7 預測引擎解碼
    (
        "讓我們解碼預測引擎的內部機制。"
        "標準化後的回歸係數揭示了真實影響力。"
        "研發支出係數：正三萬八千——壓倒性的正向驅動力。"
        "行銷支出係數：正三千五百——有補充價值。"
        "行政管理係數：負一千八百——微弱的阻力。"
        "換句話說，研發的影響力是行銷的十倍以上。"
        "每一塊投入研發的資金，都在以最高效率轉化為利潤。"
    ),
    # 8 效能驗證
    (
        "最終結果。"
        "測試集 R 平方：零點九零零一。"
        "商業目標：零點九零。門檻達到，精確到小數點後四位。"
        "訓練集與測試集差距微小，顯示模型具備良好的泛化能力，沒有過擬合。"
        "RMSE 八千九百九十六美元，MAE 六千五百九十七美元。"
        "誤差在實際投資決策的可接受範圍內。"
        "這個模型，可以被信任。"
    ),
    # 9 異常診斷
    (
        "這裡有一個重要的統計警告。"
        "我們的五折交叉驗證出現了負 R 平方，均值負零點一二，標準差高達零點六七。"
        "這不是模型崩潰，這是「小樣本極端抽取」的統計假象。"
        "五十筆資料做五折，每折只有八筆測試資料，極易抽到不具代表性的樣本。"
        "這就是為什麼單一保留測試集，十筆具代表性的樣本，"
        "在此案例中是更可靠的評估指標。"
        "了解統計假象，才能做出正確的模型選擇。"
    ),
    # 10 終極洞察
    (
        "這是整個分析最重要的洞察。"
        "很多人以為，在矽谷或紐約的黃金地段能帶來溢價。"
        "數據告訴我們：錯了。"
        "地理位置——零票，p 值超過零點九五，完全無效。"
        "而研發支出——五票全中。"
        "R&D 不是普通特徵，它是決定新創獲利能力的護城河。"
        "資本配置的核心策略只有一個：把每一分錢盡可能投入研發。"
    ),
    # 11 資本配置法則
    (
        "讓我把數據分析轉化為你可以立即使用的投資法則。"
        "黃金法則：每增加一萬美元的研發預算，預期將產生約八千美元的額外利潤。"
        "投資報酬率——研發每投入一塊錢，回收約八毛錢的年度利潤。"
        "實戰案例：研發一百五十萬、行政一百二十萬、行銷三百萬，"
        "預測年度利潤：十七萬五千八百六十美元。"
        "模型已封裝完成，可即時整合至投資組合儀表板。"
        "這不只是一份報告，這是你的資本配置決策工具。"
        "感謝各位。"
    ),
]

# 確保旁白數量匹配頁數
while len(NARRATIONS) < n_pages:
    NARRATIONS.append("請參閱投影片內容。")
NARRATIONS = NARRATIONS[:n_pages]

# ── 3. 生成中文語音（gTTS zh-TW）──────────────────────────────────────────────
print("\n" + "=" * 60)
print("步驟 2：生成創投家風格中文語音旁白")
print("=" * 60)

from gtts import gTTS

audio_files = []
for i, narr in enumerate(NARRATIONS):
    mp3_path = os.path.join(ASSETS, f"narr_{i:02d}.mp3")
    tts = gTTS(text=narr, lang="zh-tw", slow=False)
    tts.save(mp3_path)
    audio_files.append(mp3_path)
    # 取得時長
    from moviepy import AudioFileClip
    with AudioFileClip(mp3_path) as a:
        dur = a.duration
    print(f"  旁白 {i+1:02d}/{n_pages}  ({dur:.1f}s)")

# ── 4. 製作 PPTX ───────────────────────────────────────────────────────────────
print("\n" + "=" * 60)
print("步驟 3：製作 PowerPoint (PPTX)")
print("=" * 60)

from pptx import Presentation
from pptx.util import Emu

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(5143500)
blank = prs.slide_layouts[6]

for i, (png, narr) in enumerate(zip(slide_pngs, NARRATIONS)):
    sl = prs.slides.add_slide(blank)
    sl.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
    tf = sl.notes_slide.notes_text_frame
    tf.text = f"【旁白腳本】\n{narr}"

prs.save(OUT_PPTX)
print(f"[已儲存] {OUT_PPTX}  ({os.path.getsize(OUT_PPTX)/1024/1024:.1f} MB)")

# ── 5. 組合 MP4 ────────────────────────────────────────────────────────────────
print("\n" + "=" * 60)
print("步驟 4：組合 MP4 影片")
print("=" * 60)

from moviepy import AudioFileClip

TMP_DIR = os.path.join(ASSETS, "tmp_segments")
os.makedirs(TMP_DIR, exist_ok=True)

segment_files = []
for i, (png, mp3) in enumerate(zip(slide_pngs, audio_files)):
    with AudioFileClip(mp3) as audio:
        dur = audio.duration + 1.2   # 1.2s 緩衝
    seg = os.path.join(TMP_DIR, f"seg_{i:02d}.mp4")
    cmd = [
        FFMPEG, "-y",
        "-loop", "1", "-framerate", "1", "-i", png,
        "-i", mp3,
        "-c:v", "libx264", "-preset", "fast", "-crf", "23",
        "-pix_fmt", "yuv420p",
        "-vf", f"scale={TARGET_W}:{TARGET_H}",
        "-r", "24",
        "-c:a", "aac", "-b:a", "128k",
        "-t", str(dur),
        seg
    ]
    print(f"  [{i+1:02d}/{n_pages}] 編碼片段…")
    result = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="replace")
    if result.returncode != 0:
        print(f"  ERROR:\n{result.stderr[-500:]}")
        raise RuntimeError(f"ffmpeg segment {i} failed")
    segment_files.append(seg)

# concat
concat_list = os.path.join(TMP_DIR, "concat.txt")
with open(concat_list, "w", encoding="utf-8") as f:
    for seg in segment_files:
        f.write(f"file '{seg.replace(os.sep, '/')}'\n")

print(f"\n串接 {len(segment_files)} 個片段…")
tmp_out = OUT_MP4.replace(".mp4", "_tmp.mp4")
cmd = [FFMPEG, "-y", "-f", "concat", "-safe", "0", "-i", concat_list, "-c", "copy", tmp_out]
subprocess.run(cmd, capture_output=True, check=True)

# faststart
cmd = [FFMPEG, "-y", "-i", tmp_out, "-c", "copy", "-movflags", "+faststart", OUT_MP4]
result = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="replace")
if result.returncode == 0:
    os.remove(tmp_out)
else:
    os.replace(tmp_out, OUT_MP4)

size_mb = os.path.getsize(OUT_MP4) / 1024 / 1024
total_s = sum(
    AudioFileClip(f).duration + 1.2 for f in audio_files
)
print(f"\n[完成]")
print(f"  PPTX → {OUT_PPTX}")
print(f"  MP4  → {OUT_MP4}  ({size_mb:.1f} MB,  {total_s/60:.1f} 分鐘)")
