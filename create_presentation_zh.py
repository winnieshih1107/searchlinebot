"""
create_presentation_zh.py
=========================
繁體中文版：
  pres_assets_zh/slide_XX.png   — 1920x1080 投影片圖片
  pres_assets_zh/narr_XX.mp3    — gTTS 中文語音（zh-TW）
  hw6_presentation_zh.pptx      — PowerPoint（備忘稿含旁白腳本）
  hw6_presentation_zh.mp4       — 最終影片（含語音旁白）
"""
import os, textwrap, warnings
warnings.filterwarnings("ignore")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
import numpy as np

BASE   = r"D:\wi"
ASSETS = os.path.join(BASE, "pres_assets_zh")
os.makedirs(ASSETS, exist_ok=True)

W, H = 1920, 1080
DPI  = 100

# ── 中文字型：優先使用微軟正黑體 ────────────────────────────────────────────
def get_zh_font(size):
    for fname in ["Microsoft JhengHei", "Microsoft YaHei", "SimHei", "DFKai-SB"]:
        try:
            return fm.FontProperties(family=fname, size=size)
        except Exception:
            pass
    return fm.FontProperties(size=size)

ZH_FONT_FAMILIES = ["Microsoft JhengHei", "Microsoft YaHei", "SimHei",
                     "DFKai-SB", "sans-serif"]
matplotlib.rcParams["font.family"] = ZH_FONT_FAMILIES


# ─────────────────────────────────────────────────────────────────────────────
# 投影片內容（繁體中文）
# ─────────────────────────────────────────────────────────────────────────────
SLIDES = [
    # 0 ── 封面 ───────────────────────────────────────────────────────────────
    dict(
        type="title",
        title="50 個新創公司",
        sub1="專家級 CRISP-DM 迴歸分析",
        sub2="利用機器學習預測新創公司年度利潤",
        sub3="作業六  ·  2026-06-12",
        color="#2C3E50",
        narration=(
            "歡迎收看五十個新創公司的機器學習分析。"
            "本次報告依循完整的 CRISP-DM 六階段方法論，"
            "從業務理解到部署，"
            "建立一套可預測新創公司年度利潤的迴歸模型。"
        ),
    ),
    # 1 ── 執行摘要 ────────────────────────────────────────────────────────────
    dict(
        type="kv",
        title="執行摘要",
        color="#1A5276",
        rows=[
            ("資料集",        "50 間新創公司  ×  5 個欄位"),
            ("預測目標",      "年度利潤（連續數值，美元）"),
            ("最佳模型",      "Lasso 迴歸（基準特徵組合）"),
            ("交叉驗證 R²",   "0.9317  ±  0.0380（共 50 筆評分）"),
            ("交叉驗證 RMSE", "9,527 美元"),
            ("測試集 R²",     "0.9001"),
            ("測試集 RMSE",   "8,996 美元"),
            ("關鍵驅動因素",  "研發支出   Pearson r = 0.973"),
            ("6 個模型皆超標？", "是  ✓  全數超越 R² ≥ 0.90 目標"),
        ],
        narration=(
            "執行摘要顯示各項指標表現優異。"
            "最佳模型 Lasso 迴歸，使用三個基準特徵，"
            "交叉驗證 R 平方達到零點九三一七，"
            "均方根誤差僅九千五百二十七美元。"
            "測試集 R 平方為零點九零零一，超越目標。"
            "研發支出是利潤的核心驅動力，皮爾森相關係數高達零點九七三。"
            "六個模型全數通過零點九零的門檻。"
        ),
    ),
    # 2 ── CRISP-DM 流程圖 ────────────────────────────────────────────────────
    dict(
        type="image",
        title="CRISP-DM 分析流程",
        color="#1A5276",
        image=os.path.join(BASE, "hw6_workflow_excalidraw.png"),
        caption="六階段 CRISP-DM 管線，包含模型效果不佳時的重新調整回饋迴圈",
        narration=(
            "我們遵循 CRISP-DM 六個階段："
            "業務理解、資料理解、資料準備、建模、評估，以及部署。"
            "關鍵設計是評估鑽石節點後的回饋迴圈——"
            "若交叉驗證指標未達標，則退回第四階段重新調整模型。"
            "這個迭代架構確保只在效能驗證後才進行部署。"
        ),
    ),
    # 3 ── 第一階段：業務理解 ──────────────────────────────────────────────────
    dict(
        type="bullets",
        title="第一階段  —  業務理解",
        color="#2E86C1",
        items=[
            ("問題定義",    "根據新創公司的預算配置，預測其年度利潤"),
            ("任務類型",    "監督式迴歸  —  目標為連續美元數值"),
            ("輸入特徵",    "研發支出、行政費用、行銷支出、所在州別"),
            ("預測目標",    "年度利潤（連續數值，美元）"),
            ("成功標準 R²", "交叉驗證 R² ≥ 0.90"),
            ("成功標準 RMSE", "交叉驗證 RMSE < 15,000 美元"),
            ("專家洞察",    "研發支出是「公司品質」的潛在變數——捕捉創新、人才與智慧財產"),
            ("專家警告",    "n=50 時，普通 5 折交叉驗證不穩定，需改用 RepeatedKFold"),
        ],
        narration=(
            "業務目標是根據新創公司的預算配置，預測其年度利潤。"
            "這是一個監督式迴歸問題，目標為連續的美元數值。"
            "我們設定了明確的成功標準："
            "交叉驗證 R 平方至少須達零點九零，"
            "RMSE 須低於一萬五千美元。"
            "專家洞察指出，研發支出代表公司品質這個潛在變數，"
            "它涵蓋創新管線、人才質量與智慧財產，共同驅動長期獲利能力。"
        ),
    ),
    # 4 ── 第二階段：資料理解 ──────────────────────────────────────────────────
    dict(
        type="split_image",
        title="第二階段  —  資料理解",
        color="#17A589",
        left_title="關鍵統計",
        left_items=[
            "50 筆資料  ×  5 個欄位",
            "無遺漏值",
            "研發支出  r = 0.9729  ★（極強相關）",
            "行銷支出  r = 0.7478（強相關）",
            "行政費用  r = 0.2007（弱相關）",
            "州別：三組利潤範圍高度重疊（無預測力）",
        ],
        image=os.path.join(BASE, "phase2_understanding.png"),
        narration=(
            "資料集包含五十間新創公司、五個欄位，無任何遺漏值。"
            "最重要的發現是相關係數表。"
            "研發支出與利潤的皮爾森相關係數高達零點九七，"
            "單一特徵便能解釋約 95% 的利潤變異。"
            "行銷支出排名第二，相關係數為零點七五。"
            "行政費用的相關係數僅零點二零，幾乎沒有預測力。"
            "三個州別的利潤範圍高度重疊，確認州別在此樣本量下毫無預測信號。"
        ),
    ),
    # 5 ── 第三階段：資料準備 ──────────────────────────────────────────────────
    dict(
        type="bullets",
        title="第三階段  —  資料準備",
        color="#229954",
        items=[
            ("特徵工程",   "研發比例 = 研發支出 / 總支出   |   Log 研發 = log1p(研發支出)"),
            ("特徵篩選",   "5 種演算法：SFS、RFE、SelectKBest、Lasso、RF"),
            ("刪除州別",   "0/5 票  •  OLS p 值 > 0.95  •  n=50 樣本不足"),
            ("刪除行政費用", "0/5 票  •  OLS p 值 = 0.61  •  無獨立預測信號"),
            ("管線設計",   "StandardScaler → 模型  （僅在訓練集 fit，避免資料洩漏）"),
            ("訓練/測試分割", "80% 訓練  /  20% 測試  •  random_state = 42"),
            ("關鍵修正",   "KFold(5)  →  RepeatedKFold(5 折 × 10 次) = 50 筆評分"),
            ("穩定性改善", "R² 標準差：0.67  →  0.04  （提升 17 倍！）"),
        ],
        narration=(
            "在資料準備階段，我們工程化了兩個新特徵："
            "研發比例捕捉研究強度，為研發支出佔總支出的比例；"
            "Log 研發捕捉研發投資的邊際遞減效應。"
            "執行五種特徵篩選演算法後，"
            "州別和行政費用各獲得零票，一致被排除。"
            "所有模型皆包在標準化縮放管線中，防止資料洩漏。"
            "最關鍵的方法論修正是將普通五折交叉驗證"
            "替換為 5 折乘以 10 次的 RepeatedKFold，獲得 50 筆驗證分數。"
            "這使 R 平方標準差從零點六七降至零點零四，提升十七倍穩定性。"
        ),
    ),
    # 6 ── 第四階段：建模 ──────────────────────────────────────────────────────
    dict(
        type="kv",
        title="第四階段  —  建模（交叉驗證結果）",
        color="#7D3C98",
        rows=[
            ("模型",              "CV R²      標準差      RMSE"),
            ("Lasso  ★ 最佳",     "0.9317    0.038     9,527 美元"),
            ("線性迴歸",          "0.9315    0.038     9,536 美元"),
            ("ElasticNet",        "0.9313    0.038     9,549 美元"),
            ("Ridge 嶺迴歸",      "0.9309    0.038     9,578 美元"),
            ("隨機森林",          "0.9285    0.046     9,923 美元"),
            ("梯度提升",          "0.9018    0.051    11,495 美元"),
            ("─────────────", "──────────────────────────────"),
            ("成功標準",          "≥ 0.90   < 0.10   < 15,000 美元"),
            ("結果",              "✓  6 個模型全數通過"),
        ],
        narration=(
            "我們訓練了六個模型，每個皆包在標準化縮放管線中。"
            "所有模型使用五十筆交叉驗證分數進行評估。"
            "結果非常一致——四個線性模型均聚集在 R 平方零點九三一附近，"
            "標準差僅零點零三八，相當穩定。"
            "隨機森林和梯度提升方差略高，但仍超越目標。"
            "Lasso 以 R 平方零點九三一七、RMSE 九千五百二十七美元微幅勝出。"
            "六個模型全數通過三項成功標準。"
        ),
    ),
    # 7 ── 第五階段：評估（結果圖）──────────────────────────────────────────────
    dict(
        type="image",
        title="第五階段  —  評估（9 格綜合結果圖）",
        color="#CA6F1E",
        image=os.path.join(BASE, "results_figure.png"),
        caption="CV R²、穩定性、RMSE、舊vs新KFold對比、散佈圖、相關矩陣、排行榜",
        narration=(
            "九格評估圖完整呈現分析結果。"
            "第一列確認所有模型的 R 平方超標且變異低。"
            "關鍵對比圖顯示——舊的五折交叉驗證在完全相同的資料上"
            "可能產生從負零點八一到正零點八九的劇烈波動。"
            "而 RepeatedKFold 每次都穩定在零點九三附近。"
            "散佈圖確認研發支出的主導地位，OLS 的 R 平方達零點九四七。"
            "排行榜彙整了最終模型排名。"
        ),
    ),
    # 8 ── 特徵篩選 ────────────────────────────────────────────────────────────
    dict(
        type="image",
        title="特徵篩選  —  5 種方法比較",
        color="#CA6F1E",
        image=os.path.join(BASE, "feature_selection_figure.png"),
        caption="SFS / RFE / SelectKBest / Lasso / 隨機森林  ·  RMSE 與 R² 隨特徵數量的變化",
        narration=(
            "我們比較了五種特徵篩選演算法，針對六個候選特徵進行評估。"
            "結論明確：最佳特徵數 k 等於二——研發支出加行銷支出——"
            "在所有有效方法中皆達到最佳效能，RMSE 約八千八百八十四美元，R 平方零點九三八九。"
            "重要警告：遞迴特徵消除法（RFE）在此資料集表現極差。"
            "k 等於一到三時，RFE 的 R 平方為負值，"
            "因為它將州別虛擬變數排在最高優先，選出的特徵完全無用。"
            "這是 RFE 與稀疏類別特徵結合時的已知缺陷。"
        ),
    ),
    # 9 ── 第六階段：部署 ──────────────────────────────────────────────────────
    dict(
        type="bullets",
        title="第六階段  —  部署",
        color="#943126",
        items=[
            ("模型儲存",    'joblib.dump(pipeline, "50_startups_best_model.pkl")'),
            ("預測 API",    "predict_profit(rd_spend, administration, marketing_spend)"),
            ("高研發（160k）", "→  預測年度利潤  ≈  185,000 美元"),
            ("中研發（80k）",  "→  預測年度利潤  ≈  115,000 美元"),
            ("零研發",         "→  預測年度利潤  ≈   38,000 美元"),
            ("監控機制",    "每月追蹤 RMSE 對比滾動實績"),
            ("重新訓練觸發", "RMSE > 12,000 美元  或  測試 R² < 0.85"),
            ("擴展觸發",    "樣本數達 n ≥ 200 時，重新執行完整 CRISP-DM"),
        ],
        narration=(
            "部署階段，最佳管線——標準化縮放接 Lasso——已用 joblib 序列化儲存。"
            "predict_profit 函式接受研發支出、行政費用、行銷支出三個輸入，"
            "回傳預測年度利潤。"
            "三個範例預測展示模型的邏輯性："
            "高研發支出十六萬美元預測利潤約十八萬五千美元，"
            "中等研發八萬美元預測利潤約十一萬五千美元，"
            "零研發則只能預測約三萬八千美元，印證模型合理性。"
            "已建立每月 RMSE 監控，並設定明確的重新訓練觸發條件。"
        ),
    ),
    # 10 ── 專家小組討論 ────────────────────────────────────────────────────────
    dict(
        type="bullets",
        title="專家小組討論  —  5 回合",
        color="#1A5276",
        items=[
            ("第 1 回：問題定義", "研發 = 公司品質的潛在變數  •  州別：樣本量不足"),
            ("第 2 回：特徵重要性", "研發主導（r=0.97）  •  Log 轉換捕捉邊際遞減效應"),
            ("第 3 回：建模策略", "Lasso 產生稀疏解  •  n=50 必須用 RepeatedKFold"),
            ("第 4 回：共識決策", "保留：研發 + 行銷  •  刪除：行政費用 + 州別（全數同意）"),
            ("第 5 回：部署考量", "RMSE 9k～10k 對投資決策已足夠  •  每月監控"),
            ("",                   ""),
            ("最終共識", "k=2 特徵  ·  Lasso  ·  RepeatedKFold  ·  每月監控"),
        ],
        narration=(
            "我們的專家小組由行銷、研發、州政府政策和銷售四位專家組成，"
            "進行了五回合的結構化討論。"
            "第一回合一致同意：研發支出是公司品質的潛在變數。"
            "第二回合支持對研發支出進行對數轉換以捕捉邊際遞減效應。"
            "第三回合選定 Lasso 和 RepeatedKFold 作為方法論核心。"
            "第四回合達成一致共識：保留研發支出和行銷支出，刪除行政費用和州別。"
            "第五回合確認九千到一萬美元的均方根誤差，對投資決策完全可接受。"
            "最終建議：二個特徵、Lasso 正規化、RepeatedKFold、每月監控。"
        ),
    ),
    # 11 ── 關鍵發現 ───────────────────────────────────────────────────────────
    dict(
        type="bullets",
        title="關鍵發現",
        color="#1A5276",
        items=[
            ("發現 1",  "研發支出單一特徵解釋 94.7% 的利潤變異（r=0.973）"),
            ("發現 2",  "6 個模型全數超越 CV R² ≥ 0.90 目標，且有充裕空間"),
            ("發現 3",  "RepeatedKFold 關鍵——R² 標準差：0.67 → 0.04（提升 17 倍）"),
            ("發現 4",  "最佳特徵數 k=2（研發支出 + 行銷支出）"),
            ("發現 5",  "州別無用——0/5 票，OLS p 值 > 0.95"),
            ("發現 6",  "行政費用無獨立信號——OLS p 值 = 0.61"),
            ("發現 7",  "工程化特徵在 n=50 時增加不穩定性（需 n > 100）"),
            ("發現 8",  "⚠  RFE 危險：k=1 時 R² = -0.24（選中州別虛擬變數）"),
            ("發現 9",  "樹模型具競爭力，但在 n=50 時未優於線性模型"),
            ("發現 10", "最佳：Lasso  ·  CV R²=0.9317  ·  RMSE=$9,527  ·  Test R²=0.9001"),
        ],
        narration=(
            "十大關鍵發現如下。"
            "第一，研發支出單一特徵解釋近九十五趴的利潤變異。"
            "第二，六個模型全數超越 R 平方目標。"
            "第三，RepeatedKFold 對小樣本資料至關重要——將 R 平方標準差提升十七倍穩定性。"
            "第四，最佳特徵數僅二個。"
            "第五和第六，州別和行政費用在 n=50 時無法提供有效信號。"
            "第七，工程化特徵需要更多資料才能穩定。"
            "第八，這是嚴重警告——遞迴特徵消除法在此資料集上，k 等於一到三時會產生負 R 平方。"
            "第九，樹模型具競爭力但未優於線性模型。"
            "第十，最佳模型測試集 R 平方達零點九零零一，超越所有目標。"
        ),
    ),
    # 12 ── 結論 ───────────────────────────────────────────────────────────────
    dict(
        type="conclusion",
        title="結論",
        color="#2C3E50",
        points=[
            "研發投資是新創公司獲利能力的最關鍵驅動力",
            "2 個特徵的 Lasso 模型達到 R² = 0.93——對 n=50 而言極為出色",
            "RepeatedKFold 是小樣本迴歸評估的必要方法",
            "特徵篩選共識：在此資料集刪除州別與行政費用",
            "模型已部署為 50_startups_best_model.pkl，提供 predict_profit() API",
            "建議：蒐集更多資料（n ≥ 200）並重新執行完整 CRISP-DM 流程",
        ],
        narration=(
            "結論："
            "一旦找到正確的驅動因素，五十個新創公司的問題其實相當單純。"
            "研發投資——作為公司品質的代理指標——"
            "以壓倒性的優勢決定了公司獲利能力。"
            "一個只有兩個特徵的 Lasso 模型，便達到 R 平方零點九三，"
            "遠超零點九零的目標。"
            "關鍵方法論貢獻是 RepeatedKFold 交叉驗證，"
            "顯著提升了小樣本評估的穩定性。"
            "模型已完成部署，可供投資組合決策使用。"
            "感謝您收看本次五十個新創公司的專家級 CRISP-DM 分析。"
        ),
    ),
]

TOTAL = len(SLIDES)

# ─────────────────────────────────────────────────────────────────────────────
# 色彩設定
# ─────────────────────────────────────────────────────────────────────────────
LIGHT_BG = {
    "#2C3E50": "#ECF0F1", "#1A5276": "#EBF5FB",
    "#2E86C1": "#EBF5FB", "#17A589": "#E8F8F5",
    "#229954": "#EAFAF1", "#7D3C98": "#F5EEF8",
    "#CA6F1E": "#FEF9E7", "#943126": "#FDEDEC",
}


def draw_base(title, color, slide_num):
    fig = plt.figure(figsize=(W / DPI, H / DPI), dpi=DPI)
    fig.patch.set_facecolor("white")
    bg = LIGHT_BG.get(color, "#F8F9FA")

    # 標題列（上方 20%）
    ax_hdr = fig.add_axes([0, 0.80, 1, 0.20])
    ax_hdr.set_facecolor(color)
    ax_hdr.axis("off")
    ax_hdr.plot([0, 1], [0.97, 0.97], color="white", lw=5, alpha=0.25,
                transform=ax_hdr.transAxes)
    ax_hdr.text(0.04, 0.50, title,
                transform=ax_hdr.transAxes,
                fontsize=28, fontweight="bold", color="white",
                va="center", ha="left", clip_on=False)
    ax_hdr.text(0.97, 0.50, f"{slide_num}/{TOTAL}",
                transform=ax_hdr.transAxes,
                fontsize=13, color="white", alpha=0.75,
                va="center", ha="right")

    # 頁腳（下方 5.5%）
    ax_ftr = fig.add_axes([0, 0, 1, 0.055])
    ax_ftr.set_facecolor("#ECF0F1")
    ax_ftr.axis("off")
    ax_ftr.text(0.02, 0.5, "50 個新創公司  —  專家級 CRISP-DM 分析  |  作業六",
                transform=ax_ftr.transAxes,
                fontsize=11, color="#7F8C8D", va="center")

    # 內容區
    ax = fig.add_axes([0.04, 0.08, 0.92, 0.70])
    ax.set_facecolor(bg)
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    return fig, ax


# ── 各類型投影片渲染函式 ───────────────────────────────────────────────────────

def render_title(slide, idx):
    fig = plt.figure(figsize=(W / DPI, H / DPI), dpi=DPI)
    fig.patch.set_facecolor(slide["color"])
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_facecolor(slide["color"])
    ax.axis("off")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    for yy, alpha in [(0.93, 0.18), (0.88, 0.10)]:
        ax.fill_between([0, 1], [yy, yy], [yy + 0.06, yy + 0.06],
                        color="white", alpha=alpha, transform=ax.transAxes)
    ax.text(0.5, 0.63, slide["title"], transform=ax.transAxes,
            fontsize=68, fontweight="bold", color="white", va="center", ha="center")
    ax.text(0.5, 0.50, slide["sub1"], transform=ax.transAxes,
            fontsize=30, color="white", alpha=0.90, va="center", ha="center")
    ax.text(0.5, 0.40, slide["sub2"], transform=ax.transAxes,
            fontsize=19, color="white", alpha=0.75, va="center", ha="center")
    ax.text(0.5, 0.28, slide["sub3"], transform=ax.transAxes,
            fontsize=15, color="white", alpha=0.60, va="center", ha="center")
    ax_bot = fig.add_axes([0, 0, 1, 0.06])
    ax_bot.set_facecolor("#1A252F"); ax_bot.axis("off")
    ax_bot.text(0.5, 0.5, "CRISP-DM  ·  scikit-learn  ·  Python  ·  繁體中文版",
                transform=ax_bot.transAxes,
                fontsize=12, color="white", alpha=0.6, va="center", ha="center")
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


def render_kv(slide, idx):
    fig, ax = draw_base(slide["title"], slide["color"], idx + 1)
    rows = slide["rows"]
    n = len(rows)
    for i, row in enumerate(rows):
        y = 0.96 - i * (0.92 / n)
        if len(row) >= 2 and row[1].startswith("──"):
            ax.plot([0.01, 0.99], [y - 0.01, y - 0.01],
                    color="#BDC3C7", lw=1.2, transform=ax.transAxes)
            continue
        label = row[0]
        value = row[1] if len(row) > 1 else ""
        bold = i == 0 or i >= n - 2
        ax.text(0.02, y, label, transform=ax.transAxes,
                fontsize=15, fontweight="bold" if bold else "normal",
                color=slide["color"], va="top", ha="left")
        ax.text(0.38, y, value, transform=ax.transAxes,
                fontsize=14, color="#2C3E50", va="top", ha="left",
                fontweight="bold" if ("★" in value or "✓" in value) else "normal")
        if i < n - 1:
            ax.plot([0.01, 0.99], [y - 0.048, y - 0.048],
                    color="#D5D8DC", lw=0.6, transform=ax.transAxes)
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


def render_bullets(slide, idx):
    fig, ax = draw_base(slide["title"], slide["color"], idx + 1)
    items = slide["items"]
    n = len(items)
    for i, (lbl, txt) in enumerate(items):
        y = 0.97 - i * (0.91 / n)
        if not lbl and not txt:
            continue
        ax.add_patch(mpatches.FancyBboxPatch(
            (0.01, y - 0.032), 0.007, 0.038,
            boxstyle="round,pad=0.002",
            facecolor=slide["color"], edgecolor="none",
            transform=ax.transAxes, zorder=3))
        ax.text(0.025, y - 0.012, lbl, transform=ax.transAxes,
                fontsize=13, fontweight="bold", color=slide["color"],
                va="center", ha="left")
        ax.text(0.27, y - 0.012, txt, transform=ax.transAxes,
                fontsize=12, color="#2C3E50", va="center", ha="left",
                linespacing=1.3, wrap=True)
        ax.plot([0.01, 0.99], [y - 0.058, y - 0.058],
                color="#D5D8DC", lw=0.5, transform=ax.transAxes)
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


def render_image(slide, idx):
    fig, ax = draw_base(slide["title"], slide["color"], idx + 1)
    img_path = slide.get("image", "")
    if img_path and os.path.isfile(img_path):
        img = plt.imread(img_path)
        ax.imshow(img, aspect="auto",
                  extent=[0.005, 0.990, 0.04, 0.960],
                  transform=ax.transAxes, zorder=2)
    caption = slide.get("caption", "")
    if caption:
        ax.text(0.5, 0.012, caption, transform=ax.transAxes,
                fontsize=10.5, color="#7F8C8D", ha="center", va="bottom",
                style="italic")
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


def render_split_image(slide, idx):
    fig, ax = draw_base(slide["title"], slide["color"], idx + 1)
    items = slide["left_items"]
    n = len(items)
    ax.text(0.02, 0.97, slide.get("left_title", ""), transform=ax.transAxes,
            fontsize=14, fontweight="bold", color=slide["color"], va="top")
    for i, txt in enumerate(items):
        y = 0.88 - i * (0.80 / n)
        ax.add_patch(mpatches.FancyBboxPatch(
            (0.01, y - 0.025), 0.007, 0.032,
            boxstyle="round,pad=0.002",
            facecolor=slide["color"], edgecolor="none",
            transform=ax.transAxes, zorder=3))
        ax.text(0.026, y - 0.008, txt, transform=ax.transAxes,
                fontsize=12, color="#2C3E50", va="center", ha="left")
    ax.plot([0.47, 0.47], [0.02, 0.97], color=slide["color"], lw=1.5, alpha=0.3,
            transform=ax.transAxes)
    img_path = slide.get("image", "")
    if img_path and os.path.isfile(img_path):
        img = plt.imread(img_path)
        ax.imshow(img, aspect="auto",
                  extent=[0.48, 0.990, 0.02, 0.970],
                  transform=ax.transAxes, zorder=2)
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


def render_conclusion(slide, idx):
    fig, ax = draw_base(slide["title"], slide["color"], idx + 1)
    points = slide["points"]
    n = len(points)
    for i, pt in enumerate(points):
        y = 0.90 - i * (0.72 / n)
        ax.text(0.025, y, "✔", transform=ax.transAxes,
                fontsize=20, color="#27AE60", va="center")
        ax.text(0.075, y, pt, transform=ax.transAxes,
                fontsize=13.5, color="#2C3E50", va="center",
                fontweight="bold" if i == 0 else "normal")
    ax_ban = fig.add_axes([0.05, 0.09, 0.90, 0.10])
    ax_ban.set_facecolor(slide["color"]); ax_ban.axis("off")
    ax_ban.text(0.5, 0.5,
                "感謝收看  ·  GitHub: winnieshih1107/50_Startups_hw6",
                transform=ax_ban.transAxes,
                fontsize=14, color="white", va="center", ha="center")
    path = os.path.join(ASSETS, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=DPI, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return path


RENDERERS = {
    "title":       render_title,
    "kv":          render_kv,
    "bullets":     render_bullets,
    "image":       render_image,
    "split_image": render_split_image,
    "conclusion":  render_conclusion,
}


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────
def build_pptx(slide_pngs):
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    blank = prs.slide_layouts[6]
    for sd, png in zip(SLIDES, slide_pngs):
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
        tf = sl.notes_slide.notes_text_frame
        tf.text = f"【旁白腳本】\n{sd['narration']}"
    out = os.path.join(BASE, "hw6_presentation_zh.pptx")
    prs.save(out)
    print(f"[已儲存] {out}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# 主程式
# ─────────────────────────────────────────────────────────────────────────────
def main():
    slide_pngs  = []
    audio_files = []

    print("=" * 60)
    print(f"正在產生 {TOTAL} 張投影片 + 中文語音旁白…")
    print("=" * 60)

    for i, slide in enumerate(SLIDES):
        print(f"\n  投影片 {i+1:02d}/{TOTAL}  [{slide['type']}]  {slide.get('title','')[:40]}")

        # ── 渲染投影片 PNG ─────────────────────────────────────────────────
        renderer = RENDERERS.get(slide["type"], render_bullets)
        png_path = renderer(slide, i)
        slide_pngs.append(png_path)
        print(f"    PNG  -> {os.path.basename(png_path)}")

        # ── 產生中文語音 ────────────────────────────────────────────────────
        mp3_path = os.path.join(ASSETS, f"narr_{i:02d}.mp3")
        try:
            from gtts import gTTS
            tts = gTTS(text=slide["narration"], lang="zh-tw", slow=False)
            tts.save(mp3_path)
            print(f"    MP3  -> {os.path.basename(mp3_path)}")
        except Exception as e:
            print(f"    gTTS 錯誤 ({e})，使用靜音替代")
            import wave
            wav_path = mp3_path.replace(".mp3", ".wav")
            with wave.open(wav_path, "w") as wf:
                wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(22050)
                wf.writeframes(b"\x00\x00" * 22050 * 4)
            mp3_path = wav_path
        audio_files.append(mp3_path)

    # ── PPTX ──────────────────────────────────────────────────────────────────
    print("\n" + "=" * 60)
    print("正在建立 PPTX…")
    build_pptx(slide_pngs)

    # ── MP4 ───────────────────────────────────────────────────────────────────
    print("正在合成影片（moviepy）…")
    try:
        from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

        clips = []
        for png_path, mp3_path in zip(slide_pngs, audio_files):
            audio    = AudioFileClip(mp3_path)
            duration = audio.duration + 1.0
            clip     = (ImageClip(png_path)
                        .with_duration(duration)
                        .with_audio(audio))
            clips.append(clip)

        final   = concatenate_videoclips(clips, method="compose")
        mp4_out = os.path.join(BASE, "hw6_presentation_zh.mp4")

        final.write_videofile(
            mp4_out,
            fps=24,
            codec="libx264",
            audio_codec="aac",
            ffmpeg_params=["-crf", "23", "-preset", "fast",
                           "-movflags", "+faststart"],
            logger=None,
        )
        total_s = sum(AudioFileClip(f).duration for f in audio_files)
        print(f"[已儲存] {mp4_out}")
        print(f"         總時長：{total_s/60:.1f} 分鐘")

    except Exception as e:
        import traceback
        print(f"  MP4 錯誤：{e}")
        traceback.print_exc()

    print("\n" + "=" * 60)
    print("完成！")
    print(f"  PPTX  -> {os.path.join(BASE, 'hw6_presentation_zh.pptx')}")
    print(f"  MP4   -> {os.path.join(BASE, 'hw6_presentation_zh.mp4')}")
    print("=" * 60)


if __name__ == "__main__":
    main()
